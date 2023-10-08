from pptx import Presentation


def pptx2text(file_path):
    """
    Extract text content from a PowerPoint (.pptx) file.

    Parameters:
    - file_path (str): The path to the PowerPoint file.

    Returns:
    - str: The extracted text content.

    Example:
    ```python
    text_content = pptx2text("path/to/presentation.pptx")
    print(text_content)
    ```

    :param file_path: The path to the PowerPoint file.
    :type file_path: str
    :return: The extracted text content.
    :rtype: str
    """
    # Load the PowerPoint presentation
    presentation = Presentation(file_path)

    # Initialize an empty list to store extracted texts
    texts = []

    # Iterate through slides and shapes to extract text content
    for slide in presentation.slides:
        for shape in slide.shapes:
            # Check if the shape has a text attribute
            if hasattr(shape, "text"):
                # Exclude placeholder text "‹#›"
                if shape.text != "‹#›":
                    # Process and append text content, handling newline characters
                    texts.append(' '.join([
                        word.strip() for word in shape.text.split('\n')
                    ]))

                    try:
                        # Clean up and punctuate the last added text
                        texts[-1] = texts[-1].strip()
                        if texts[-1][-1] not in ['?', '!', '@', '"', "'", '.', '\t']:
                            texts[-1] += '.'
                    except Exception as e:
                        pass

    # Combine all extracted texts into a single string
    return ' '.join([text.strip() for text in texts])

# Uncomment the following line to test the function
# print(pptx2text("backend/src/pptx2pdf/test.pptx"))
