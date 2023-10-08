from pptx import Presentation


def pptx2text(file_path):
    f = Presentation(file_path)
    texts = []
    for slide in f.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if (shape.text != "‹#›"):
                    texts.append(shape.text)

    return texts
