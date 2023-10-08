from pptx import Presentation


def pptx2text(file_path):
    f = Presentation(file_path)
    texts = []
    for slide in f.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if (shape.text != "‹#›"):
                    texts.append(' '.join([
                        word.strip() for word in shape.text.split('\n')
                    ]))
                try:
                    texts[-1] = texts[-1].strip()
                    if (texts[-1][-1] not in ['?', '!', '@', '"', "'", '.', '\t']):
                        texts[-1] += '.'
                except Exception as e:
                    pass

    return ' '.join(
        [text.strip() for text in texts]
    )


# pptx2text("backend/src/pptx2pdf/test.pptx")
# print(pptx2text("backend/src/pptx2pdf/test.pptx"))
